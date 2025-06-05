from fastapi import FastAPI, File, UploadFile, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
import os
import shutil
from pptx import Presentation
from typing import List, Tuple
import uuid

app = FastAPI()

# Allow CORS (adjust for your frontend domain in production)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

UPLOAD_DIR = "uploads"
RESULT_DIR = "results"
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(RESULT_DIR, exist_ok=True)


def extract_text_from_slide(slide) -> List[str]:
    text_items = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text_items.append(shape.text)
    return text_items


def dummy_grammar_correction(text: str) -> str:
    # Replace this with real grammar/narrative correction logic or API
    return text.replace("bad", "good").replace("mistkae", "mistake")


def process_pptx(file_path: str) -> Tuple[str, int, List[dict]]:
    prs = Presentation(file_path)
    amended_count = 0
    corrections = []

    for idx, slide in enumerate(prs.slides):
        original_texts = extract_text_from_slide(slide)
        corrected_texts = [dummy_grammar_correction(text) for text in original_texts]

        if original_texts != corrected_texts:
            amended_count += 1

        # Replace text in slides
        for shape, corrected in zip(
            [s for s in slide.shapes if hasattr(s, "text")], corrected_texts
        ):
            shape.text = corrected

        corrections.append({
            "slide_number": idx + 1,
            "original": original_texts,
            "corrected": corrected_texts,
        })

    corrected_filename = f"corrected_{uuid.uuid4().hex}.pptx"
    corrected_path = os.path.join(RESULT_DIR, corrected_filename)
    prs.save(corrected_path)

    return corrected_path, amended_count, corrections


@app.post("/api/process-ppt")
async def process_ppt(file: UploadFile = File(...), report: str = Form(...), options: str = Form(...)):
    if not file.filename.endswith(".pptx"):
        return JSONResponse(status_code=400, content={"error": "Only .pptx files are supported."})

    file_id = uuid.uuid4().hex
    uploaded_path = os.path.join(UPLOAD_DIR, f"{file_id}.pptx")

    with open(uploaded_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    corrected_path, amended_slides, corrections = process_pptx(uploaded_path)

    download_url = f"/api/download/{os.path.basename(corrected_path)}"

    return {
        "amendedSlidesCount": amended_slides,
        "fileUrl": download_url,
        "corrections": corrections,
    }


@app.get("/api/download/{filename}")
def download_file(filename: str):
    path = os.path.join(RESULT_DIR, filename)
    if not os.path.exists(path):
        return JSONResponse(status_code=404, content={"error": "File not found."})
    return FileResponse(path, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename=filename)
